"""PPTX parser — `python-pptx`.

PPTX is the easiest format: every slide already has a structural skeleton
(title placeholder + body shapes). We emit one `SLIDE_TITLE` block per
slide that has a title placeholder, plus one `SLIDE_BULLET` per bullet
paragraph and one `SLIDE_BODY` block per non-bulleted paragraph.

`detect_typographic_structure` will trivially flag this as structured.
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Iterable

from agent.models import Block, BlockKind, Document, SlideLocator, SourceType

logger = logging.getLogger(__name__)


def parse(path: str | Path) -> Document:
    try:
        from pptx import Presentation  # type: ignore
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore  # noqa: F401
    except ImportError as e:
        raise RuntimeError(
            "python-pptx is not installed. Install with `pip install python-pptx`."
        ) from e

    p = Path(path)
    prs = Presentation(str(p))
    blocks = list(_iter_blocks(prs))
    if not blocks:
        raise ValueError(f"No text content found in {p}.")
    return Document(
        source_path=str(p),
        source_type=SourceType.PPTX,
        blocks=blocks,
        size_metric={"slides": len(prs.slides)},
        parser_meta={"extractor": "python-pptx"},
    )


def _iter_blocks(prs) -> Iterable[Block]:
    for slide_idx, slide in enumerate(prs.slides, start=1):
        title_text = _slide_title(slide)
        if title_text:
            yield Block(
                kind=BlockKind.SLIDE_TITLE,
                text=title_text,
                locator=SlideLocator(
                    slide_index_start=slide_idx, slide_index_end=slide_idx
                ),
            )
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if _is_title_placeholder(shape):
                continue  # already emitted as SLIDE_TITLE
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if not text:
                    continue
                kind = (
                    BlockKind.SLIDE_BULLET
                    if _is_bulleted(para)
                    else BlockKind.SLIDE_BODY
                )
                yield Block(
                    kind=kind,
                    text=text,
                    locator=SlideLocator(
                        slide_index_start=slide_idx, slide_index_end=slide_idx
                    ),
                    meta={"indent_level": getattr(para, "level", 0) or 0},
                )


def _slide_title(slide) -> str:
    if slide.shapes.title is None:
        return ""
    title = slide.shapes.title.text_frame.text.strip() if slide.shapes.title.has_text_frame else ""
    return title


def _is_title_placeholder(shape) -> bool:
    if not getattr(shape, "is_placeholder", False):
        return False
    ph = shape.placeholder_format
    if ph is None:
        return False
    # idx 0 is the title placeholder in standard PPTX layouts.
    return ph.idx == 0


def _is_bulleted(paragraph) -> bool:
    # python-pptx exposes the underlying XML; presence of <a:buChar> or
    # <a:buAutoNum> in the paragraph properties indicates a bullet.
    pPr = paragraph._pPr  # private but stable in python-pptx
    if pPr is None:
        return False
    for tag in ("buChar", "buAutoNum"):
        if pPr.find(f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag}") is not None:
            return True
    return False
